from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import textwrap

# ── colours ──────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1F, 0x38, 0x64)
BLUE      = RGBColor(0x2E, 0x75, 0xB6)
LBLUE     = RGBColor(0xBD, 0xD7, 0xEE)
RED       = RGBColor(0xC0, 0x00, 0x00)
RED_BG    = RGBColor(0xFF, 0xCC, 0xCC)
GREEN     = RGBColor(0x37, 0x56, 0x23)
GREEN_BG  = RGBColor(0xCC, 0xFF, 0xCC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY     = RGBColor(0xF2, 0xF2, 0xF2)
DGRAY     = RGBColor(0x44, 0x44, 0x44)
AMBER     = RGBColor(0xFF, 0xBF, 0x00)
AMBER_BG  = RGBColor(0xFF, 0xF0, 0xCC)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _rgb(r, g, b):
    return RGBColor(r, g, b)


def _gap_color(val):
    """Return (text_color, bg_color) based on gap value."""
    if val is None:
        return DGRAY, LGRAY
    if val > 0:
        return GREEN, GREEN_BG
    if val < 0:
        return RED, RED_BG
    return DGRAY, LGRAY


def _wow_label(val):
    if val is None or val == 0:
        return '—'
    arrow = '▲' if val > 0 else '▼'
    sign  = '+' if val > 0 else ''
    num   = str(int(round(val)))
    return f'{arrow} {sign}{num}'


def _wow_color(val):
    if val is None or val == 0:
        return DGRAY
    return GREEN if val > 0 else RED


def _strip_negative_prefix(value):
    if isinstance(value, str) and value.startswith('-'):
        return value[1:].strip()
    return value


def _positive_raw_for_metric(metric: str, value):
    if metric == 'RD' and value is not None:
        try:
            return abs(float(value))
        except (TypeError, ValueError):
            return value
    return value


def create_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = line_width
    else:
        line.fill.background()
    return shape


def _add_text_box(slide, x, y, w, h, text, font_size=Pt(10), bold=False,
                  color=None, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def _cell(slide, x, y, w, h, label, value, val_raw=None,
          label_size=Pt(7), value_size=Pt(16), is_gap=False, is_wow=False):
    """Draw a labelled KPI cell."""
    fill = LGRAY
    txt_color = DGRAY
    val_color = NAVY

    if is_gap and val_raw is not None:
        txt_color, fill = _gap_color(val_raw)
        val_color = txt_color
    if is_wow:
        val_color = _wow_color(val_raw)
        fill = LGRAY

    _add_rect(slide, x, y, w, h, fill_rgb=fill,
              line_rgb=RGBColor(0xCC, 0xCC, 0xCC))

    # label
    _add_text_box(slide, x + Inches(0.05), y + Inches(0.04),
                  w - Inches(0.1), Inches(0.22),
                  label, font_size=label_size, bold=False, color=DGRAY,
                  align=PP_ALIGN.CENTER)
    # value
    _add_text_box(slide, x + Inches(0.05), y + Inches(0.26),
                  w - Inches(0.1), h - Inches(0.3),
                  value, font_size=value_size, bold=True, color=val_color,
                  align=PP_ALIGN.CENTER)


def add_title_slide(prs, week_label, quarter_label):
    slide = _blank_slide(prs)

    # full background
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=NAVY)

    # accent stripe
    _add_rect(slide, 0, Inches(3.5), SLIDE_W, Inches(0.06), fill_rgb=BLUE)

    _add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
                  'Wipro Weekly Performance Report',
                  font_size=Pt(32), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.6),
                  f'{week_label}  |  {quarter_label}',
                  font_size=Pt(20), bold=False, color=LBLUE, align=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
                  'Resource Utilisation  ·  Resource Deallocation  ·  Net Addition',
                  font_size=Pt(13), bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1.5), Inches(6.6), Inches(10), Inches(0.4),
                  'Confidential — Internal Use Only',
                  font_size=Pt(9), bold=False, color=RGBColor(0x80, 0x80, 0x80),
                  align=PP_ALIGN.CENTER)


def add_account_metric_slide(prs, rec: dict, week_label: str, quarter_label: str,
                              chart_type: str, bar_path: str = None,
                              trend_path: str = None):
    slide = _blank_slide(prs)
    metric  = rec['metric']
    account = rec['account']

    # ── HEADER BAR ────────────────────────────────────────────────────────────
    hdr_h = Inches(0.88)
    _add_rect(slide, 0, 0, SLIDE_W, hdr_h, fill_rgb=NAVY)
    _add_rect(slide, 0, hdr_h, SLIDE_W, SLIDE_H - hdr_h, fill_rgb=RGBColor(0xF6, 0xF8, 0xFB), line_rgb=None)

    # metric badge
    badge_w = Inches(1.3)
    _add_rect(slide, Inches(0.15), Inches(0.12), badge_w, Inches(0.58),
              fill_rgb=BLUE, line_rgb=None)
    _add_text_box(slide, Inches(0.15), Inches(0.12), badge_w, Inches(0.58),
                  metric, font_size=Pt(26), bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1.6), Inches(0.16), Inches(8), Inches(0.55),
                  account, font_size=Pt(24), bold=True, color=WHITE,
                  align=PP_ALIGN.LEFT)

    _add_text_box(slide, Inches(10.5), Inches(0.2), Inches(2.7), Inches(0.45),
                  f'{week_label}  |  {quarter_label}',
                  font_size=Pt(12), bold=False, color=LBLUE,
                  align=PP_ALIGN.RIGHT)

    # ── KPI TABLE SECTION ─────────────────────────────────────────────────────
    row_y = Inches(1.08)
    row_h = Inches(0.72)
    pad   = Inches(0.15)
    col_w = (SLIDE_W - 2 * pad) / 5

    # column headers row
    hdrs = ['Plan QTR', 'WK Plan', 'WK Act', 'Gap', 'WoW']
    for i, hdr in enumerate(hdrs):
        cx = pad + i * col_w
        _add_rect(slide, cx, row_y, col_w - Inches(0.04), Inches(0.26),
                  fill_rgb=NAVY, line_rgb=None)
        _add_text_box(slide, cx + Inches(0.05), row_y + Inches(0.03),
                      col_w - Inches(0.14), Inches(0.22),
                      hdr, font_size=Pt(11), bold=True, color=WHITE,
                      align=PP_ALIGN.CENTER)

    # main metric values row
    vals     = [rec['fmt_plan_qtr'], rec['fmt_wk_plan'], rec['fmt_wk_act'],
                rec['fmt_gap'],      rec.get('fmt_wow', '—')]
    raw_vals = [rec['plan_qtr'],     rec['wk_plan'],     rec['wk_act'],
                rec['gap'],          rec.get('wow')]

    if metric == 'RD':
        vals = [(_strip_negative_prefix(v) if isinstance(v, str) else v) for v in vals]
        raw_vals = [
            raw_vals[0],
            _positive_raw_for_metric(metric, raw_vals[1]),
            _positive_raw_for_metric(metric, raw_vals[2]),
            _positive_raw_for_metric(metric, raw_vals[3]),
            raw_vals[4],
        ]

    val_y = row_y + Inches(0.27)
    val_h = Inches(0.62)

    for i, (v, rv) in enumerate(zip(vals, raw_vals)):
        cx = pad + i * col_w
        is_gap = (i == 3)
        is_wow = (i == 4)

        fill_c  = LGRAY
        val_col = NAVY

        if is_gap and rv is not None:
            _, fill_c = _gap_color(rv)
            val_col, _ = _gap_color(rv)
        if is_wow:
            val_col = _wow_color(rv)
            v = _wow_label(rv)

        _add_rect(slide, cx, val_y, col_w - Inches(0.04), val_h,
                  fill_rgb=fill_c, line_rgb=RGBColor(0xCC, 0xCC, 0xCC))
        _add_text_box(slide, cx + Inches(0.05), val_y + Inches(0.08),
                      col_w - Inches(0.14), val_h - Inches(0.1),
                      v, font_size=Pt(22), bold=True, color=val_col,
                      align=PP_ALIGN.CENTER)

    # ── BPM ROW ───────────────────────────────────────────────────────────────
    bpm_y = val_y + val_h + Inches(0.10)
    bpm_h = Inches(0.62)

    bpm_hdr_h = Inches(0.24)
    bpm_val_h = bpm_h - bpm_hdr_h

    bpm_hdrs = ['BPM Plan QTR', 'BPM WK Plan', 'BPM WK Act', 'BPM Gap', 'BPM WoW']
    bpm_vals = [rec['fmt_bpm_plan_qtr'], rec['fmt_bpm_wk_plan'], rec['fmt_bpm_wk_act'],
                rec['fmt_bpm_gap'],      rec.get('fmt_bpm_wow', '—')]
    bpm_raws = [rec['bpm_plan_qtr'],     rec['bpm_wk_plan'],     rec['bpm_wk_act'],
                rec['bpm_gap'],          rec.get('bpm_wow')]

    if metric == 'RD':
        bpm_vals = [(_strip_negative_prefix(v) if isinstance(v, str) else v) for v in bpm_vals]
        bpm_raws = [
            bpm_raws[0],
            _positive_raw_for_metric(metric, bpm_raws[1]),
            _positive_raw_for_metric(metric, bpm_raws[2]),
            _positive_raw_for_metric(metric, bpm_raws[3]),
            bpm_raws[4],
        ]

    for i, (hdr, v, rv) in enumerate(zip(bpm_hdrs, bpm_vals, bpm_raws)):
        cx = pad + i * col_w
        is_gap = (i == 3)
        is_wow = (i == 4)

        _add_rect(slide, cx, bpm_y, col_w - Inches(0.04), bpm_hdr_h,
                  fill_rgb=BLUE, line_rgb=None)
        _add_text_box(slide, cx + Inches(0.05), bpm_y + Inches(0.02),
                      col_w - Inches(0.14), bpm_hdr_h - Inches(0.04),
                      hdr, font_size=Pt(9), bold=True, color=WHITE,
                      align=PP_ALIGN.CENTER)

        fill_c  = LGRAY
        val_col = NAVY
        if is_gap and rv is not None:
            _, fill_c = _gap_color(rv)
            val_col, _ = _gap_color(rv)
        if is_wow:
            val_col = _wow_color(rv)
            v = _wow_label(rv)

        _add_rect(slide, cx, bpm_y + bpm_hdr_h, col_w - Inches(0.04), bpm_val_h,
                  fill_rgb=fill_c, line_rgb=RGBColor(0xCC, 0xCC, 0xCC))
        _add_text_box(slide, cx + Inches(0.05), bpm_y + bpm_hdr_h + Inches(0.04),
                      col_w - Inches(0.14), bpm_val_h - Inches(0.04),
                      v, font_size=Pt(16), bold=True, color=val_col,
                      align=PP_ALIGN.CENTER)

    # ── REASON SECTION ────────────────────────────────────────────────────────
    reason_y = bpm_y + bpm_h + Inches(0.14)
    reason_h = Inches(0.5)
    rp_y     = reason_y + reason_h + Inches(0.06)
    rp_h     = Inches(0.5)

    # Delta reason
    _add_rect(slide, Inches(0.12), reason_y, Inches(1.3), reason_h,
              fill_rgb=BLUE, line_rgb=None)
    _add_text_box(slide, Inches(0.12), reason_y + Inches(0.08),
                  Inches(1.3), reason_h - Inches(0.1),
                  'Delta\nReason', font_size=Pt(7.5), bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)

    dr_text = rec['delta_reason'] or '—'
    _add_rect(slide, Inches(1.45), reason_y, SLIDE_W - Inches(1.6), reason_h,
              fill_rgb=RGBColor(0xFF, 0xF8, 0xE7), line_rgb=RGBColor(0xDD, 0xDD, 0xDD))
    _add_text_box(slide, Inches(1.55), reason_y + Inches(0.05),
                  SLIDE_W - Inches(1.8), reason_h - Inches(0.1),
                  dr_text, font_size=Pt(8.5), bold=False, color=DGRAY,
                  align=PP_ALIGN.LEFT, wrap=True)

    # Recovery plan
    _add_rect(slide, Inches(0.12), rp_y, Inches(1.3), rp_h,
              fill_rgb=RGBColor(0x37, 0x56, 0x23), line_rgb=None)
    _add_text_box(slide, Inches(0.12), rp_y + Inches(0.08),
                  Inches(1.3), rp_h - Inches(0.1),
                  'Recovery\nPlan', font_size=Pt(7.5), bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)

    rp_text = rec['recovery_plan'] or '—'
    _add_rect(slide, Inches(1.45), rp_y, SLIDE_W - Inches(1.6), rp_h,
              fill_rgb=RGBColor(0xF0, 0xFF, 0xF0), line_rgb=RGBColor(0xDD, 0xDD, 0xDD))
    _add_text_box(slide, Inches(1.55), rp_y + Inches(0.05),
                  SLIDE_W - Inches(1.8), rp_h - Inches(0.1),
                  rp_text, font_size=Pt(8.5), bold=False, color=DGRAY,
                  align=PP_ALIGN.LEFT, wrap=True)

    # ── BOTTOM SECTION — chart + KPI cards ────────────────────────────────────
    bottom_y = rp_y + rp_h + Inches(0.12)
    bottom_h = SLIDE_H - bottom_y - Inches(0.1)
    half_w   = SLIDE_W / 2

    if chart_type in ('bar', 'all') and bar_path:
        slide.shapes.add_picture(bar_path,
                                 Inches(0.1), bottom_y,
                                 half_w - Inches(0.2), bottom_h)

    if chart_type in ('kpi', 'all'):
        kpi_x    = half_w + Inches(0.1) if chart_type == 'all' else Inches(0.5)
        kpi_area_w = (half_w - Inches(0.3)) if chart_type == 'all' else (SLIDE_W - Inches(1))
        kpi_items = [
            ('Gap',      rec['fmt_gap'],      rec['gap'],      True,  False),
            ('BPM Gap',  rec['fmt_bpm_gap'],  rec['bpm_gap'],  True,  False),
            ('WoW',      _wow_label(rec.get('wow')),     rec.get('wow'),     False, True),
            ('BPM WoW',  _wow_label(rec.get('bpm_wow')), rec.get('bpm_wow'), False, True),
        ]
        card_w  = (kpi_area_w - Inches(0.1)) / 2
        card_h  = (bottom_h - Inches(0.1)) / 2
        for idx, (lbl, val, raw, ig, iw) in enumerate(kpi_items):
            col = idx % 2
            row = idx // 2
            cx = kpi_x + col * (card_w + Inches(0.1))
            cy = bottom_y + row * (card_h + Inches(0.08))

            if ig and raw is not None:
                txt_c, fill_c = _gap_color(raw)
            elif iw:
                fill_c = LGRAY
                txt_c  = _wow_color(raw)
            else:
                fill_c = LGRAY
                txt_c  = NAVY

            _add_rect(slide, cx, cy, card_w, card_h,
                      fill_rgb=fill_c, line_rgb=RGBColor(0xCC, 0xCC, 0xCC))

            _add_text_box(slide, cx + Inches(0.1), cy + Inches(0.1),
                          card_w - Inches(0.2), Inches(0.28),
                          lbl, font_size=Pt(10), bold=True, color=DGRAY,
                          align=PP_ALIGN.CENTER)

            _add_text_box(slide, cx + Inches(0.05), cy + Inches(0.35),
                          card_w - Inches(0.1), card_h - Inches(0.45),
                          val, font_size=Pt(26), bold=True, color=txt_c,
                          align=PP_ALIGN.CENTER)

    # trend chart if provided
    if trend_path and chart_type == 'all':
        pass  # reserved for future line chart mode


def add_metric_overview_slide(prs, records: list, metric: str,
                               week_label: str, quarter_label: str,
                               chart_path: str = None,
                               title_prefix: str | None = None):
    """One slide showing all accounts for a single metric."""
    slide = _blank_slide(prs)

    # subtle background
    hdr_h = Inches(0.78)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=RGBColor(0xF6, 0xF8, 0xFB), line_rgb=None)
    _add_rect(slide, 0, 0, SLIDE_W, hdr_h, fill_rgb=NAVY)
    title_text = f'{metric} — {title_prefix}' if title_prefix else f'{metric} — All Accounts'
    _add_text_box(slide, Inches(0.2), Inches(0.12), Inches(9), Inches(0.55),
                  title_text,
                  font_size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _add_text_box(slide, Inches(10.5), Inches(0.18), Inches(2.7), Inches(0.42),
                  f'{week_label}  |  {quarter_label}',
                  font_size=Pt(10), bold=False, color=LBLUE, align=PP_ALIGN.RIGHT)

    # table
    col_hdrs = ['Account', 'Plan QTR', 'WK Plan', 'WK Act', 'Gap', 'WoW',
                'BPM Plan', 'BPM WK Act', 'BPM Gap', 'BPM WoW']
    col_keys = ['account', 'fmt_plan_qtr', 'fmt_wk_plan', 'fmt_wk_act',
                'fmt_gap', 'wow', 'fmt_bpm_plan_qtr', 'fmt_bpm_wk_act',
                'fmt_bpm_gap', 'bpm_wow']
    gap_cols = {4: 'gap', 8: 'bpm_gap'}
    wow_cols  = {5: 'wow', 9: 'bpm_wow'}

    tbl_y = Inches(0.92)
    tbl_h = SLIDE_H - tbl_y - Inches(0.14)

    if chart_path:
        # table on left 55%, chart on right 45%
        tbl_w = SLIDE_W * 0.55
        slide.shapes.add_picture(chart_path,
                                 tbl_w + Inches(0.1), tbl_y,
                                 SLIDE_W - tbl_w - Inches(0.2), tbl_h)
    else:
        tbl_w = SLIDE_W - Inches(0.2)

    n_cols  = len(col_hdrs)
    acct_w  = Inches(1.9)
    rest_w  = (tbl_w - Inches(0.1) - acct_w) / (n_cols - 1)

    row_h      = min(Inches(0.38), (tbl_h - Inches(0.38)) / max(len(records), 1))
    hdr_row_h  = Inches(0.38)

    # header row
    cx = Inches(0.1)
    for ci, hdr in enumerate(col_hdrs):
        cw = acct_w if ci == 0 else rest_w
        _add_rect(slide, cx, tbl_y, cw - Inches(0.02), hdr_row_h,
                  fill_rgb=NAVY, line_rgb=None)
        _add_text_box(slide, cx + Inches(0.03), tbl_y + Inches(0.06),
                      cw - Inches(0.08), hdr_row_h - Inches(0.1),
                      hdr, font_size=Pt(8), bold=True, color=WHITE,
                      align=PP_ALIGN.CENTER)
        cx += cw

    # data rows
    for ri, rec in enumerate(records):
        ry = tbl_y + hdr_row_h + ri * row_h
        row_fill = WHITE if ri % 2 == 0 else LGRAY
        cx = Inches(0.1)
        for ci, key in enumerate(col_keys):
            cw = acct_w if ci == 0 else rest_w
            raw_v = rec.get(gap_cols[ci]) if ci in gap_cols else (
                    rec.get(wow_cols[ci]) if ci in wow_cols else None)
            if metric == 'RD' and ci in gap_cols and raw_v is not None:
                raw_v = abs(raw_v)

            if ci in gap_cols and raw_v is not None:
                _, fill_c = _gap_color(raw_v)
                txt_c, _ = _gap_color(raw_v)
            elif ci in wow_cols:
                fill_c = row_fill
                txt_c  = _wow_color(raw_v)
            else:
                fill_c = row_fill
                txt_c  = DGRAY if ci == 0 else NAVY

            val_str = rec.get(key, '—')
            if metric == 'RD' and isinstance(val_str, str) and key.startswith('fmt_'):
                val_str = _strip_negative_prefix(val_str)
            if ci in wow_cols:
                val_str = _wow_label(raw_v)
            if val_str is None:
                val_str = '—'

            _add_rect(slide, cx, ry, cw - Inches(0.02), row_h,
                      fill_rgb=fill_c, line_rgb=RGBColor(0xDD, 0xDD, 0xDD))
            _add_text_box(slide, cx + Inches(0.04), ry + Inches(0.04),
                          cw - Inches(0.1), row_h - Inches(0.06),
                          str(val_str), font_size=Pt(8), bold=(ci == 0),
                          color=txt_c, align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
            cx += cw


def add_summary_slide(prs, all_data: dict, week_label: str, quarter_label: str):
    """Cover summary: one row per metric showing totals."""
    slide = _blank_slide(prs)

    hdr_h = Inches(0.75)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=RGBColor(0xF6, 0xF8, 0xFB), line_rgb=None)
    _add_rect(slide, 0, 0, SLIDE_W, hdr_h, fill_rgb=NAVY)
    _add_text_box(slide, Inches(0.2), Inches(0.12), Inches(9), Inches(0.55),
                  'Summary — All Metrics',
                  font_size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _add_text_box(slide, Inches(10.5), Inches(0.18), Inches(2.7), Inches(0.42),
                  f'{week_label}  |  {quarter_label}',
                  font_size=Pt(10), bold=False, color=LBLUE, align=PP_ALIGN.RIGHT)

    metrics = list(all_data.keys())
    card_y  = Inches(1.1)
    card_h  = (SLIDE_H - card_y - Inches(0.2)) / max(len(metrics), 1)

    for mi, metric in enumerate(metrics):
        records = all_data[metric]
        if not records:
            continue
        my = card_y + mi * card_h

        _add_rect(slide, Inches(0.1), my, Inches(1.5), card_h - Inches(0.08),
                  fill_rgb=BLUE, line_rgb=None)
        _add_text_box(slide, Inches(0.1), my + (card_h - Inches(0.08)) / 2 - Inches(0.2),
                      Inches(1.5), Inches(0.4),
                      metric, font_size=Pt(20), bold=True, color=WHITE,
                      align=PP_ALIGN.CENTER)

        # aggregate totals
        def _sum(key):
            vals = [r.get(key) for r in records if r.get(key) is not None]
            if metric == 'RD':
                vals = [abs(v) for v in vals if isinstance(v, (int, float))]
            return sum(vals) if vals else None

        def _fmt_s(v):
            if v is None:
                return '—'
            return str(int(round(v)))

        def _first_text(key):
            for r in records:
                val = r.get(key)
                if val not in (None, '', '—'):
                    return str(val)
            return '—'

        stats = [
            ('Accounts',   str(len(records)),     None),
            ('Plan QTR',   _first_text('plan_qtr'), None),
            ('WK Act',     _fmt_s(_sum('wk_act')),   _sum('wk_act')),
            ('Gap',        _fmt_s(_sum('gap')),       _sum('gap')),
            ('BPM Gap',    _fmt_s(_sum('bpm_gap')),   _sum('bpm_gap')),
        ]

        sw = (SLIDE_W - Inches(1.8)) / len(stats)
        for si, (lbl, val, raw) in enumerate(stats):
            sx = Inches(1.7) + si * sw
            is_gap = lbl in ('Gap', 'BPM Gap')
            if is_gap and raw is not None:
                txt_c, fill_c = _gap_color(raw)
            else:
                fill_c, txt_c = LGRAY, NAVY

            _add_rect(slide, sx, my, sw - Inches(0.08), card_h - Inches(0.08),
                      fill_rgb=fill_c, line_rgb=RGBColor(0xCC, 0xCC, 0xCC))
            _add_text_box(slide, sx + Inches(0.05), my + Inches(0.06),
                          sw - Inches(0.18), Inches(0.22),
                          lbl, font_size=Pt(9), bold=True, color=DGRAY,
                          align=PP_ALIGN.CENTER)
            _add_text_box(slide, sx + Inches(0.05), my + Inches(0.28),
                          sw - Inches(0.18), card_h - Inches(0.38),
                          val, font_size=Pt(24), bold=True, color=txt_c,
                          align=PP_ALIGN.CENTER)


def add_grouped_accounts_slide(prs, records: list, metric: str,
                                week_label: str, quarter_label: str,
                                chart_path: str = None,
                                group_name: str = ''):
    """
    One slide showing multiple tagged accounts in a side-by-side comparison table.
    Each account gets a row; columns show key metrics with colour-coded gaps.
    An optional bar chart comparing gaps sits in the bottom-right.
    """
    slide = _blank_slide(prs)
    n = len(records)

    # ── header ────────────────────────────────────────────────────────────────
    hdr_h = Inches(0.78)
    _add_rect(slide, 0, 0, SLIDE_W, hdr_h, fill_rgb=NAVY)

    account_names = ' · '.join(r['account'] for r in records)
    display_names = account_names if len(account_names) <= 90 else account_names[:87] + '…'
    title_label   = group_name if group_name else 'Comparison Slide'

    badge_w = Inches(1.3)
    _add_rect(slide, Inches(0.15), Inches(0.1), badge_w, Inches(0.58),
              fill_rgb=BLUE, line_rgb=None)
    _add_text_box(slide, Inches(0.15), Inches(0.1), badge_w, Inches(0.58),
                  metric, font_size=Pt(22), bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)

    # group name (tag) in large text
    _add_text_box(slide, Inches(1.6), Inches(0.08), Inches(8.8), Inches(0.35),
                  title_label, font_size=Pt(14), bold=True,
                  color=WHITE, align=PP_ALIGN.LEFT)
    # account names in smaller subtext
    _add_text_box(slide, Inches(1.6), Inches(0.42), Inches(8.8), Inches(0.28),
                  display_names, font_size=Pt(8.5), bold=False,
                  color=LBLUE, align=PP_ALIGN.LEFT)
    _add_text_box(slide, Inches(11.0), Inches(0.2), Inches(2.2), Inches(0.42),
                  f'{week_label}  |  {quarter_label}',
                  font_size=Pt(10), bold=False, color=LBLUE, align=PP_ALIGN.RIGHT)

    # ── layout split: table left (55%), chart right (45%) ────────────────────
    tbl_w    = SLIDE_W * (1.0 if not chart_path else 0.56)
    tbl_x    = Inches(0.1)
    tbl_y    = hdr_h + Inches(0.08)
    tbl_h    = SLIDE_H - tbl_y - Inches(0.1)

    col_defs = [
        ('Account',      Inches(1.7),  'account',        None,       False, False),
        ('Plan QTR',     Inches(0.9),  'fmt_plan_qtr',   None,       False, False),
        ('WK Plan',      Inches(0.85), 'fmt_wk_plan',    None,       False, False),
        ('WK Act',       Inches(0.85), 'fmt_wk_act',     None,       False, False),
        ('Gap',          Inches(0.85), 'fmt_gap',        'gap',      True,  False),
        ('WoW',          Inches(0.75), 'fmt_wow',        'wow',      False, True),
        ('BPM Plan',     Inches(0.85), 'fmt_bpm_plan_qtr', None,     False, False),
        ('BPM Act',      Inches(0.85), 'fmt_bpm_wk_act', None,       False, False),
        ('BPM Gap',      Inches(0.85), 'fmt_bpm_gap',    'bpm_gap',  True,  False),
        ('BPM WoW',      Inches(0.75), 'fmt_bpm_wow',    'bpm_wow',  False, True),
        ('Delta Reason', Inches(2.2),  'delta_reason',   None,       False, False),
    ]

    # trim columns to fit table width
    total_col_w = sum(c[1] for c in col_defs)
    if total_col_w > tbl_w - Inches(0.1):
        # drop Delta Reason if tight
        col_defs = col_defs[:-1]
        total_col_w = sum(c[1] for c in col_defs)

    hdr_row_h = Inches(0.36)
    data_row_h = min(Inches(0.52), (tbl_h - hdr_row_h) / max(n, 1))

    # header row
    cx = tbl_x
    for (lbl, cw, _, __, _g, _w) in col_defs:
        _add_rect(slide, cx, tbl_y, cw - Inches(0.02), hdr_row_h,
                  fill_rgb=NAVY, line_rgb=None)
        _add_text_box(slide, cx + Inches(0.04), tbl_y + Inches(0.07),
                      cw - Inches(0.1), hdr_row_h - Inches(0.1),
                      lbl, font_size=Pt(7.5), bold=True, color=WHITE,
                      align=PP_ALIGN.CENTER)
        cx += cw

    # data rows
    for ri, rec in enumerate(records):
        ry = tbl_y + hdr_row_h + ri * data_row_h
        row_fill = WHITE if ri % 2 == 0 else LGRAY
        cx = tbl_x

        for ci, (lbl, cw, key, raw_key, is_gap, is_wow) in enumerate(col_defs):
            raw_v = rec.get(raw_key) if raw_key else None
            if metric == 'RD' and is_gap and raw_v is not None:
                raw_v = abs(raw_v)

            val   = rec.get(key, '—') or '—'
            if metric == 'RD' and isinstance(val, str) and key.startswith('fmt_'):
                val = _strip_negative_prefix(val)

            if is_gap and raw_v is not None:
                txt_c, fill_c = _gap_color(raw_v)
            elif is_wow:
                fill_c = row_fill
                txt_c  = _wow_color(raw_v)
                val    = _wow_label(raw_v)
            elif ci == 0:
                fill_c = row_fill
                txt_c  = NAVY
            else:
                fill_c = row_fill
                txt_c  = DGRAY

            if isinstance(val, float):
                val = f'{val:.1f}' if val != int(val) else str(int(val))
            val = str(val)[:45]  # truncate long text

            _add_rect(slide, cx, ry, cw - Inches(0.02), data_row_h,
                      fill_rgb=fill_c, line_rgb=RGBColor(0xDD, 0xDD, 0xDD))
            align = PP_ALIGN.LEFT if ci in (0, len(col_defs) - 1) else PP_ALIGN.CENTER
            _add_text_box(slide, cx + Inches(0.04), ry + Inches(0.06),
                          cw - Inches(0.1), data_row_h - Inches(0.08),
                          val, font_size=Pt(8 if ci == 0 else 9),
                          bold=(ci == 0), color=txt_c, align=align, wrap=True)
            cx += cw

    # ── chart (right panel) ───────────────────────────────────────────────────
    if chart_path:
        chart_x = tbl_w + Inches(0.15)
        chart_w = SLIDE_W - chart_x - Inches(0.1)
        slide.shapes.add_picture(chart_path, chart_x, tbl_y, chart_w, tbl_h)
